"""Document text extraction for various file formats."""
from io import BytesIO

# Max characters to extract (enough for AI analysis, not too expensive)
MAX_EXTRACT_CHARS = 5000


def extract_text(content: bytes, mime_type: str) -> str:
    """Extract text from document based on MIME type.
    
    Returns extracted text (up to MAX_EXTRACT_CHARS) or empty string on failure.
    """
    extractors = {
        "application/pdf": _extract_pdf,
        "application/msword": _extract_doc_fallback,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
        "application/vnd.ms-excel": _extract_xlsx_fallback,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_xlsx,
        "application/vnd.ms-powerpoint": _extract_pptx_fallback,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
        "text/plain": _extract_txt,
    }
    
    extractor = extractors.get(mime_type)
    if not extractor:
        return ""
    
    try:
        text = extractor(content)
        return _clean_text(text)[:MAX_EXTRACT_CHARS]
    except Exception:
        return ""


def _clean_text(text: str) -> str:
    """Clean extracted text - remove excessive whitespace."""
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    return "\n".join(lines)


def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF using PyMuPDF."""
    import fitz  # PyMuPDF
    
    doc = fitz.open(stream=content, filetype="pdf")
    texts = []
    
    for page in doc:
        texts.append(page.get_text())
        if len("\n".join(texts)) > MAX_EXTRACT_CHARS:
            break
    
    doc.close()
    return "\n".join(texts)


def _extract_docx(content: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    
    doc = Document(BytesIO(content))
    texts = []
    
    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)
        if len("\n".join(texts)) > MAX_EXTRACT_CHARS:
            break
    
    return "\n".join(texts)


def _extract_xlsx(content: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook
    
    wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
    texts = []
    
    for sheet in wb.sheetnames[:3]:  # Max 3 sheets
        ws = wb[sheet]
        texts.append(f"[{sheet}]")
        
        for row in ws.iter_rows(max_row=50, values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text:
                texts.append(row_text)
        
        if len("\n".join(texts)) > MAX_EXTRACT_CHARS:
            break
    
    wb.close()
    return "\n".join(texts)


def _extract_pptx(content: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    
    prs = Presentation(BytesIO(content))
    texts = []
    
    for i, slide in enumerate(prs.slides[:20], 1):  # Max 20 slides
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        
        if slide_texts:
            texts.append(f"[Slide {i}] " + " ".join(slide_texts))
        
        if len("\n".join(texts)) > MAX_EXTRACT_CHARS:
            break
    
    return "\n".join(texts)


def _extract_txt(content: bytes) -> str:
    """Extract text from plain text file."""
    for encoding in ["utf-8", "cp1255", "iso-8859-8", "latin-1"]:
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return ""


def _extract_doc_fallback(content: bytes) -> str:
    """Fallback for old .doc format - limited support."""
    # Old .doc format is binary, best effort extraction
    try:
        text = content.decode("utf-8", errors="ignore")
        # Filter printable characters
        return "".join(c for c in text if c.isprintable() or c in "\n\t")
    except Exception:
        return ""


def _extract_xlsx_fallback(content: bytes) -> str:
    """Fallback for old .xls format."""
    return ""  # Old Excel format requires xlrd, skip for now


def _extract_pptx_fallback(content: bytes) -> str:
    """Fallback for old .ppt format."""
    return ""  # Old PowerPoint format not supported
